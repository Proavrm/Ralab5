"""
api/qualite.py — Endpoints Qualité complets
"""
from __future__ import annotations
from collections import Counter
from datetime import datetime
import io
import json
import secrets
from pathlib import Path
from typing import Optional, Any
from pydantic import BaseModel
from fastapi import APIRouter, HTTPException, Query, Depends, UploadFile, File
from app.models.qualite import (
    EquipmentCreateSchema, EquipmentUpdateSchema, EquipmentResponseSchema,
    MetrologyCreateSchema, MetrologyUpdateSchema, MetrologyResponseSchema,
    ProcedureCreateSchema, ProcedureUpdateSchema, ProcedureResponseSchema,
    StandardCreateSchema, StandardUpdateSchema, StandardResponseSchema,
    NcCreateSchema, NcUpdateSchema, NcResponseSchema,
    CATEGORIES_EQ, STATUTS_EQ, CONTROL_TYPES, CONTROL_STATUTS,
    PROC_FAMILIES, STD_FAMILIES, DOC_STATUTS, NC_SOURCES, NC_SEVERITES, NC_STATUTS,
)
from app.repositories.qualite_repository import (
    EquipmentRepository, MetrologyRepository,
    ProcedureRepository, StandardRepository, NcRepository,
    get_stats,
)
from app.core.database import connect_qsse_db
from app.services.qsse_import_service import QsseImportService, WorkbookSource
from app.services.qsse_rex_draft_service import QsseRexDraftService

router = APIRouter()
PROJECT_ROOT = Path(__file__).resolve().parents[3]
QSSE_STORAGE_ROOT = PROJECT_ROOT / "storage" / "qsse"
QSSE_ATTACHMENTS_ROOT = QSSE_STORAGE_ROOT / "fnc"
QSSE_ATTACHMENT_EXTENSIONS = {".pdf"}
QSSE_ATTACHMENT_SIZE_LIMIT = 25 * 1024 * 1024
QSSE_ATTACHMENT_REGISTER_CODES = {"FNC", "PASD", "BP", "FAE"}
QSSE_ANALYSIS_ATTACHMENTS_ROOT = QSSE_STORAGE_ROOT / "analysis"
QSSE_ANALYSIS_FNC_CODE = "fnc_global"
QSSE_ANALYSIS_PPTX_EXTENSIONS = {".pptx"}
QSSE_ANALYSIS_TEMPLATE_FILES = {
    "exploitation": QSSE_STORAGE_ROOT / "2026.02.25_Réunion exploitation.pptx",
    "codir": QSSE_STORAGE_ROOT / "ORGANIGRAMME 2026.pptx",
}
QSSE_ANALYSIS_TEMPLATE_MODES = set(QSSE_ANALYSIS_TEMPLATE_FILES.keys())
QSSE_EDITABLE_COLUMN_FIELDS = {
    "agency",
    "site",
    "theme",
    "title",
    "cause",
    "treatment",
    "status",
    "severity",
    "date_event",
    "date_closed",
    "date_saisie",
    "person",
    "pilot",
    "document_reference",
    "amount_value",
}

_eq    = EquipmentRepository()
_metro = MetrologyRepository()
_proc  = ProcedureRepository()
_std   = StandardRepository()
_nc    = NcRepository()
_qsse_rex = QsseRexDraftService()

def _resp_eq(r)   -> EquipmentResponseSchema:   return EquipmentResponseSchema(**r.__dict__ if hasattr(r,'__dict__') else {f: getattr(r,f) for f in r.__dataclass_fields__})
def _resp_m(r)    -> MetrologyResponseSchema:   return MetrologyResponseSchema(**{f: getattr(r,f) for f in r.__dataclass_fields__})
def _resp_proc(r) -> ProcedureResponseSchema:   return ProcedureResponseSchema(**{f: getattr(r,f) for f in r.__dataclass_fields__})
def _resp_std(r)  -> StandardResponseSchema:    return StandardResponseSchema(**{f: getattr(r,f) for f in r.__dataclass_fields__})
def _resp_nc(r)   -> NcResponseSchema:          return NcResponseSchema(**{f: getattr(r,f) for f in r.__dataclass_fields__})

def _to_dict(record):
    return {f: getattr(record, f) for f in record.__dataclass_fields__}


def _qsse_record_row(conn, record_id: int):
    return conn.execute(
        """
        SELECT id, register_code, record_kind, title
        FROM qsse_records
        WHERE id = ?
        """,
        (int(record_id),),
    ).fetchone()


def _serialize_qsse_document(record_id: int, row) -> dict[str, Any]:
    stored_name = str(row["stored_name"] or "")
    suffix = Path(stored_name).suffix.lower()
    return {
        "id": int(row["id"]),
        "record_id": int(record_id),
        "stored_name": stored_name,
        "original_name": row["original_name"] or stored_name,
        "content_type": row["content_type"] or "application/octet-stream",
        "file_size": int(row["file_size"] or 0),
        "created_at": row["created_at"] or "",
        "extension": suffix,
        "is_pdf": suffix == ".pdf",
        "url": f"/api/storage/qsse/fnc/{int(record_id)}/{stored_name}",
    }


def _serialize_qsse_analysis_document(row) -> dict[str, Any]:
    analysis_code = str(row["analysis_code"] or QSSE_ANALYSIS_FNC_CODE).strip() or QSSE_ANALYSIS_FNC_CODE
    source_year = int(row["source_year"] or 0)
    stored_name = str(row["stored_name"] or "")
    return {
        "id": int(row["id"]),
        "analysis_code": analysis_code,
        "source_year": source_year,
        "stored_name": stored_name,
        "original_name": row["original_name"] or stored_name,
        "content_type": row["content_type"] or "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "file_size": int(row["file_size"] or 0),
        "created_at": row["created_at"] or "",
        "url": f"/api/storage/qsse/analysis/{analysis_code}/{source_year}/{stored_name}",
    }


def _qsse_bucket_counter(rows: list[Any], key: str, limit: int = 8) -> list[tuple[str, int]]:
    counter: Counter[str] = Counter()
    for row in rows:
        label = str(row[key] or "").strip() or "Sans valeur"
        counter[label] += 1
    return counter.most_common(limit)


def _clear_presentation_slides(prs) -> None:
    # Remove all template content while preserving theme/layout masters.
    slide_ids = list(prs.slides._sldIdLst)
    for slide_id in slide_ids:
        rel_id = slide_id.rId
        prs.part.drop_rel(rel_id)
        prs.slides._sldIdLst.remove(slide_id)


def _add_bullet_lines(text_frame, lines: list[str]) -> None:
    text_frame.clear()
    for index, line in enumerate(lines):
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        paragraph.text = line


def _set_slide_shape_text(prs, slide_index: int, shape_index: int, value: str) -> None:
    if slide_index < 0 or slide_index >= len(prs.slides):
        return
    slide = prs.slides[slide_index]
    if shape_index < 0 or shape_index >= len(slide.shapes):
        return
    shape = slide.shapes[shape_index]
    if not getattr(shape, "has_text_frame", False):
        return
    shape.text_frame.clear()
    shape.text_frame.paragraphs[0].text = value


def _apply_fnc_template_content(prs, rows: list[Any], year: Optional[int], template_mode: str) -> bool:
    if len(prs.slides) < 8:
        return False

    total = len(rows)
    closed_count = sum(1 for row in rows if str(row["date_closed"] or "").strip())
    open_count = total - closed_count
    cost_sum = sum(float(row["amount_value"] or 0.0) for row in rows)
    closure_rate = round((closed_count / total) * 100.0, 1) if total else 0.0
    top_agencies = _qsse_bucket_counter(rows, "agency", limit=3)
    top_status = _qsse_bucket_counter(rows, "status", limit=3)

    agencies_text = " / ".join(f"{label}: {count}" for label, count in top_agencies) or "Aucune agence"
    status_text = " / ".join(f"{label}: {count}" for label, count in top_status) or "Aucun statut"

    latest_rows = rows[:5]
    latest_text = "\n".join(
        f"- FNC-{int(row['id'])} | {str(row['agency'] or '').strip() or 'Sans agence'} | {(str(row['title'] or '').strip() or 'Sans sujet')[:90]}"
        for row in latest_rows
    ) or "- Aucun événement FNC trouvé"

    perimeter = str(year) if year else "toutes années"
    title_line = f"Réunion {'CODIR' if template_mode == 'codir' else 'Exploitation'} - Analyse FNC {perimeter}"
    kpi_block = (
        f"Synthèse FNC {perimeter}\n"
        f"- Total: {total}\n"
        f"- Ouvertes: {open_count}\n"
        f"- Clôturées: {closed_count} ({closure_rate}%)\n"
        f"- Coût déclaré: {int(round(cost_sum))} EUR\n"
        f"- Top agences: {agencies_text}\n"
        f"- Top statuts: {status_text}"
    )

    _set_slide_shape_text(prs, 0, 2, title_line)
    _set_slide_shape_text(prs, 2, 4, "Objectifs QSSE FNC - pilotage data")
    _set_slide_shape_text(prs, 2, 9, kpi_block)
    _set_slide_shape_text(prs, 3, 4, "Situation FNC actuelle")
    _set_slide_shape_text(
        prs,
        3,
        7,
        (
            f"FNC remontées: {total}\n"
            f"Clôture: {closure_rate}%\n"
            f"Coût consolidé: {int(round(cost_sum))} EUR\n\n"
            f"Top agences: {agencies_text}\n"
            f"Top statuts: {status_text}"
        ),
    )
    _set_slide_shape_text(prs, 12, 4, "Plan d'actions CODIR / exploitation")
    _set_slide_shape_text(
        prs,
        12,
        7,
        (
            "- Prioriser les FNC majeures encore ouvertes\n"
            "- Suivre les plans d'actions hebdomadaires en agence\n"
            "- Renforcer le partage REX sur les causes récurrentes\n"
            "- Actualiser la revue FNC à chaque import"
        ),
    )
    _set_slide_shape_text(prs, 14, 3, "Suivi indicateurs FNC")
    _set_slide_shape_text(prs, 14, 5, f"FNC {perimeter} - total {total}")
    # Slide 14 has no large content zone — shape[7] is the page-number box.
    # Add a text box in the main content area (below the subtitle badge).
    if 14 < len(prs.slides):
        try:
            from pptx.util import Emu, Pt
            slide14 = prs.slides[14]
            txb = slide14.shapes.add_textbox(Emu(1157938), Emu(1280000), Emu(8800000), Emu(4500000))
            tf = txb.text_frame
            tf.word_wrap = True
            _add_bullet_lines(tf, latest_text.split("\n"))
            for para in tf.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(11)
        except Exception:
            pass
    return True


def _resolve_qsse_analysis_template(mode: str) -> Path | None:
    template_path = QSSE_ANALYSIS_TEMPLATE_FILES.get(mode)
    if template_path and template_path.exists():
        return template_path

    fallback = QSSE_ANALYSIS_TEMPLATE_FILES.get("exploitation")
    if fallback and fallback.exists():
        return fallback
    return None


def _build_fnc_analysis_pptx(rows: list[Any], year: Optional[int], template_mode: str) -> bytes:
    try:
        from pptx import Presentation
    except Exception as exc:
        raise HTTPException(
            status_code=500,
            detail="La génération PPTX nécessite python-pptx côté backend.",
        ) from exc

    template_path = _resolve_qsse_analysis_template(template_mode)
    if template_path is not None:
        prs = Presentation(str(template_path))
        if len(prs.slides) < 8 and template_mode == "codir":
            fallback = QSSE_ANALYSIS_TEMPLATE_FILES.get("exploitation")
            if fallback and fallback.exists():
                prs = Presentation(str(fallback))
        if _apply_fnc_template_content(prs, rows, year, template_mode):
            out = io.BytesIO()
            prs.save(out)
            return out.getvalue()
        _clear_presentation_slides(prs)
    else:
        prs = Presentation()

    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = "Analyse globale FNC"
    subtitle = title_slide.placeholders[1]
    subtitle.text = (
        f"Version: {template_mode.upper()}\n"
        f"Périmètre: {year if year else 'toutes années'}\n"
        f"Généré le {datetime.now().strftime('%d/%m/%Y %H:%M')}"
    )

    total = len(rows)
    closed_count = sum(1 for row in rows if str(row["date_closed"] or "").strip())
    open_count = total - closed_count
    cost_sum = sum(float(row["amount_value"] or 0.0) for row in rows)
    closure_rate = round((closed_count / total) * 100.0, 1) if total else 0.0

    kpi_slide = prs.slides.add_slide(prs.slide_layouts[1])
    kpi_slide.shapes.title.text = "Indicateurs clés"
    text_frame = kpi_slide.shapes.placeholders[1].text_frame
    _add_bullet_lines(text_frame, [
        f"Total FNC: {total}",
        f"FNC ouvertes: {open_count}",
        f"FNC clôturées: {closed_count}",
        f"Taux de clôture: {closure_rate}%",
        f"Coût déclaré: {int(round(cost_sum))} EUR",
    ])

    status_slide = prs.slides.add_slide(prs.slide_layouts[1])
    status_slide.shapes.title.text = "Répartitions principales"
    status_frame = status_slide.shapes.placeholders[1].text_frame
    status_lines = ["Top statuts"]
    for label, count in _qsse_bucket_counter(rows, "status", limit=6):
        status_lines.append(f"- {label}: {count}")
    status_lines.append("")
    status_lines.append("Top sévérités")
    for label, count in _qsse_bucket_counter(rows, "severity", limit=6):
        status_lines.append(f"- {label}: {count}")
    _add_bullet_lines(status_frame, status_lines)

    agency_slide = prs.slides.add_slide(prs.slide_layouts[1])
    agency_slide.shapes.title.text = "Agences et sujets FNC"
    agency_frame = agency_slide.shapes.placeholders[1].text_frame
    agency_lines = ["Top agences"]
    for label, count in _qsse_bucket_counter(rows, "agency", limit=8):
        agency_lines.append(f"- {label}: {count}")
    agency_lines.append("")
    agency_lines.append("Dernières FNC")
    for row in rows[:8]:
        title = str(row["title"] or "").strip() or "Sans sujet"
        reference = f"FNC-{int(row['id'])}"
        agency = str(row["agency"] or "").strip() or "Sans agence"
        agency_lines.append(f"- {reference} | {agency} | {title[:90]}")
    _add_bullet_lines(agency_frame, agency_lines)

    next_actions_slide = prs.slides.add_slide(prs.slide_layouts[1])
    next_actions_slide.shapes.title.text = "Actions proposées"
    actions_frame = next_actions_slide.shapes.placeholders[1].text_frame
    _add_bullet_lines(actions_frame, [
        "1. Prioriser les FNC majeures encore ouvertes.",
        "2. Suivre un point hebdomadaire de clôture avec les agences les plus exposées.",
        "3. Partager un retour REX ciblé sur les 5 sujets les plus fréquents.",
        "4. Mettre à jour cette présentation après chaque import QSSE.",
    ])

    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()

class EquipmentOptionSchema(BaseModel):
    value: str
    label: str
    equipment_id: int
    code: str
    equipment_label: str
    category: str
    domain: Optional[str] = None
    status: str
    serial_number: Optional[str] = None
    last_metrology: Optional[str] = None
    next_metrology: Optional[str] = None
    calibration_date: Optional[str] = None


class QsseRecordCellUpdateSchema(BaseModel):
    mode: str = "column"
    field: str
    value: Any = None
    aliases: list[str] = []


def _qsse_json_object(value: Any) -> dict[str, Any]:
    if not value:
        return {}
    if isinstance(value, dict):
        return value
    try:
        parsed = json.loads(value)
    except Exception:
        return {}
    return parsed if isinstance(parsed, dict) else {}


def _qsse_text_value(value: Any) -> str:
    return str(value or "").strip()


def _qsse_number_value(value: Any) -> float | None:
    if value is None:
        return None
    raw = str(value).strip().replace(" ", "")
    if not raw:
        return None
    try:
        return float(raw.replace(",", "."))
    except ValueError as exc:
        raise HTTPException(status_code=400, detail="Valeur numérique invalide.") from exc


EQUIPMENT_USAGE_DEFAULTS: dict[str, dict[str, Any]] = {
    "gammadensimetre": {
        "category": "Terrain",
        "status": "En service",
        "terms": ["gamma", "gammadens", "densim", "pqi", "troxler", "nucléaire", "nucleaire"],
    },
    "gammadensimetre_de": {
        "category": "Terrain",
        "status": "En service",
        "terms": ["gamma", "gammadens", "densim", "pqi", "troxler", "nucléaire", "nucleaire"],
    },
    "sondage_carotte_sc": {
        "category": "Terrain",
        "status": "En service",
        "terms": ["carotte", "carotier", "carottage", "foreuse", "couronne", "sondage", "sondeuse", "perceuse"],
    },
}


def _normalize_text(value: Optional[str]) -> str:
    return (value or "").strip().lower()


def _equipment_matches_usage(record, usage: Optional[str]) -> bool:
    if not usage:
        return True

    usage_key = usage.strip().lower()
    config = EQUIPMENT_USAGE_DEFAULTS.get(usage_key)
    if not config:
        return True

    terms = config.get("terms") or []
    if not terms:
        return True

    searchable = " ".join(
        [
            _normalize_text(record.code),
            _normalize_text(record.label),
            _normalize_text(record.domain),
            _normalize_text(record.serial_number),
            _normalize_text(record.notes),
        ]
    )

    return any(term in searchable for term in terms)


def _equipment_option(record) -> EquipmentOptionSchema:
    code = (record.code or "").strip()
    label = (record.label or "").strip()
    serial = (record.serial_number or "").strip()

    if code and label and serial:
        option_label = f"{code} - {label} ({serial})"
    elif code and label:
        option_label = f"{code} - {label}"
    else:
        option_label = label or code or str(record.uid)

    value = code or label or str(record.uid)
    calibration_date = record.last_metrology or ""

    return EquipmentOptionSchema(
        value=value,
        label=option_label,
        equipment_id=record.uid,
        code=code,
        equipment_label=label,
        category=record.category,
        domain=record.domain,
        status=record.status,
        serial_number=record.serial_number,
        last_metrology=record.last_metrology,
        next_metrology=record.next_metrology,
        calibration_date=calibration_date,
    )
# ── Stats ─────────────────────────────────────────────────────────────────────
@router.get("/stats")
def qualite_stats():
    return get_stats()


# ── Meta (listes de valeurs) ──────────────────────────────────────────────────
@router.get("/meta")
def qualite_meta():
    return {
        "categories_eq": CATEGORIES_EQ, "statuts_eq": STATUTS_EQ,
        "control_types": CONTROL_TYPES, "control_statuts": CONTROL_STATUTS,
        "proc_families": PROC_FAMILIES, "std_families": STD_FAMILIES,
        "doc_statuts": DOC_STATUTS, "nc_sources": NC_SOURCES,
        "nc_severites": NC_SEVERITES, "nc_statuts": NC_STATUTS,
    }


# ── Équipements ───────────────────────────────────────────────────────────────
@router.get("/equipment", response_model=list[EquipmentResponseSchema])
def list_equipment(
    search:   Optional[str] = Query(None),
    category: Optional[str] = Query(None),
    status:   Optional[str] = Query(None),
    labo_code: Optional[str] = Query(None),
):
    return [EquipmentResponseSchema(**_to_dict(r)) for r in _eq.all(search, category, status, labo_code)]


@router.get("/equipment-options", response_model=list[EquipmentOptionSchema])
def list_equipment_options(
    usage: Optional[str] = Query(None),
    search: Optional[str] = Query(None),
    category: Optional[str] = Query(None),
    status: Optional[str] = Query(None),
    include_inactive: bool = Query(False),
):
    usage_key = usage.strip().lower() if usage else None
    defaults = EQUIPMENT_USAGE_DEFAULTS.get(usage_key or "", {})

    resolved_category = category or defaults.get("category")
    resolved_status = status or defaults.get("status")

    if include_inactive:
        resolved_status = status

    rows = _eq.all(
        search=search,
        category=resolved_category,
        status=resolved_status,
    )

    filtered = [row for row in rows if _equipment_matches_usage(row, usage_key)]

    return [_equipment_option(row) for row in filtered]


@router.get("/equipment/{uid}", response_model=EquipmentResponseSchema)
def get_equipment(uid: int):
    r = _eq.get(uid)
    if not r: raise HTTPException(404, "Équipement non trouvé")
    return EquipmentResponseSchema(**_to_dict(r))


@router.post("/equipment", response_model=EquipmentResponseSchema, status_code=201)
def create_equipment(data: EquipmentCreateSchema):
    return EquipmentResponseSchema(**_to_dict(_eq.create(data)))


@router.put("/equipment/{uid}", response_model=EquipmentResponseSchema)
def update_equipment(uid: int, data: EquipmentUpdateSchema):
    r = _eq.update(uid, data)
    if not r: raise HTTPException(404, "Équipement non trouvé")
    return EquipmentResponseSchema(**_to_dict(r))


@router.delete("/equipment/{uid}", status_code=204)
def delete_equipment(uid: int):
    _eq.delete(uid)


# ── Métrologie par équipement ─────────────────────────────────────────────────
@router.get("/equipment/{uid}/metrology", response_model=list[MetrologyResponseSchema])
def list_metrology_for_eq(uid: int):
    return [MetrologyResponseSchema(**_to_dict(r)) for r in _metro.for_equipment(uid)]


@router.get("/metrology/alerts", response_model=list[MetrologyResponseSchema])
def metrology_alerts(days: int = Query(60)):
    return [MetrologyResponseSchema(**_to_dict(r)) for r in _metro.alerts(days)]


@router.get("/metrology/{uid}", response_model=MetrologyResponseSchema)
def get_metrology(uid: int):
    r = _metro.get(uid)
    if not r: raise HTTPException(404)
    return MetrologyResponseSchema(**_to_dict(r))


@router.post("/equipment/{eq_uid}/metrology", response_model=MetrologyResponseSchema, status_code=201)
def create_metrology(eq_uid: int, data: MetrologyCreateSchema):
    data.equipment_id = eq_uid
    return MetrologyResponseSchema(**_to_dict(_metro.create(data)))


@router.put("/metrology/{uid}", response_model=MetrologyResponseSchema)
def update_metrology(uid: int, data: MetrologyUpdateSchema):
    r = _metro.update(uid, data)
    if not r: raise HTTPException(404)
    return MetrologyResponseSchema(**_to_dict(r))


@router.delete("/metrology/{uid}", status_code=204)
def delete_metrology(uid: int):
    _metro.delete(uid)


# ── Procédures ────────────────────────────────────────────────────────────────
@router.get("/procedures", response_model=list[ProcedureResponseSchema])
def list_procedures(
    search: Optional[str] = Query(None),
    family: Optional[str] = Query(None),
    status: Optional[str] = Query(None),
):
    return [ProcedureResponseSchema(**_to_dict(r)) for r in _proc.all(search, family, status)]


@router.get("/procedures/{uid}", response_model=ProcedureResponseSchema)
def get_procedure(uid: int):
    r = _proc.get(uid)
    if not r: raise HTTPException(404)
    return ProcedureResponseSchema(**_to_dict(r))


@router.post("/procedures", response_model=ProcedureResponseSchema, status_code=201)
def create_procedure(data: ProcedureCreateSchema):
    return ProcedureResponseSchema(**_to_dict(_proc.create(data)))


@router.put("/procedures/{uid}", response_model=ProcedureResponseSchema)
def update_procedure(uid: int, data: ProcedureUpdateSchema):
    r = _proc.update(uid, data)
    if not r: raise HTTPException(404)
    return ProcedureResponseSchema(**_to_dict(r))


@router.delete("/procedures/{uid}", status_code=204)
def delete_procedure(uid: int):
    _proc.delete(uid)


# ── Normes ────────────────────────────────────────────────────────────────────
@router.get("/standards", response_model=list[StandardResponseSchema])
def list_standards(
    search: Optional[str] = Query(None),
    family: Optional[str] = Query(None),
    status: Optional[str] = Query(None),
):
    return [StandardResponseSchema(**_to_dict(r)) for r in _std.all(search, family, status)]


@router.get("/standards/{uid}", response_model=StandardResponseSchema)
def get_standard(uid: int):
    r = _std.get(uid)
    if not r: raise HTTPException(404)
    return StandardResponseSchema(**_to_dict(r))


@router.post("/standards", response_model=StandardResponseSchema, status_code=201)
def create_standard(data: StandardCreateSchema):
    return StandardResponseSchema(**_to_dict(_std.create(data)))


@router.put("/standards/{uid}", response_model=StandardResponseSchema)
def update_standard(uid: int, data: StandardUpdateSchema):
    r = _std.update(uid, data)
    if not r: raise HTTPException(404)
    return StandardResponseSchema(**_to_dict(r))


@router.delete("/standards/{uid}", status_code=204)
def delete_standard(uid: int):
    _std.delete(uid)


# ── Non-conformités ───────────────────────────────────────────────────────────
@router.get("/nc", response_model=list[NcResponseSchema])
def list_nc(
    search:      Optional[str] = Query(None),
    status:      Optional[str] = Query(None),
    severity:    Optional[str] = Query(None),
    source_type: Optional[str] = Query(None),
):
    return [NcResponseSchema(**_to_dict(r)) for r in _nc.all(search, status, severity, source_type)]


@router.get("/nc/{uid}", response_model=NcResponseSchema)
def get_nc(uid: int):
    r = _nc.get(uid)
    if not r: raise HTTPException(404)
    return NcResponseSchema(**_to_dict(r))


@router.post("/nc", response_model=NcResponseSchema, status_code=201)
def create_nc(data: NcCreateSchema):
    return NcResponseSchema(**_to_dict(_nc.create(data)))


@router.put("/nc/{uid}", response_model=NcResponseSchema)
def update_nc(uid: int, data: NcUpdateSchema):
    r = _nc.update(uid, data)
    if not r: raise HTTPException(404)
    return NcResponseSchema(**_to_dict(r))


@router.delete("/nc/{uid}", status_code=204)
def delete_nc(uid: int):
    _nc.delete(uid)


# ── QSSE / FNC ───────────────────────────────────────────────────────────────
@router.get("/qsse/overview")
def qsse_overview(year: Optional[int] = Query(None)):
    where_clauses = ["1 = 1"]
    params: list[Any] = []

    if year is not None:
        where_clauses.append("source_year = ?")
        params.append(int(year))

    where_sql = " AND ".join(where_clauses)

    with connect_qsse_db() as conn:
        totals = conn.execute(
            f"""
            SELECT
                COUNT(*) AS total_records,
                SUM(CASE WHEN record_kind = 'event' THEN 1 ELSE 0 END) AS event_records,
                SUM(CASE WHEN record_kind = 'indicator' THEN 1 ELSE 0 END) AS indicator_records,
                SUM(CASE WHEN source_year = 2025 THEN 1 ELSE 0 END) AS records_2025,
                SUM(CASE WHEN source_year = 2026 THEN 1 ELSE 0 END) AS records_2026
            FROM qsse_records
            WHERE {where_sql}
            """,
            params,
        ).fetchone()

        registers = conn.execute(
            f"""
            SELECT register_code, COUNT(*) AS row_count
            FROM qsse_records
            WHERE {where_sql}
            GROUP BY register_code
            ORDER BY row_count DESC, register_code ASC
            """,
            params,
        ).fetchall()

        years = conn.execute(
            """
            SELECT source_year, COUNT(*) AS row_count
            FROM qsse_records
            GROUP BY source_year
            ORDER BY source_year ASC
            """
        ).fetchall()

        latest_live_run = conn.execute(
            """
            SELECT id, source_file, source_year, source_mode, sheet_count, row_count, inserted_count, skipped_count, status, updated_at
            FROM qsse_import_runs
            WHERE source_mode = 'live'
            ORDER BY id DESC
            LIMIT 1
            """
        ).fetchone()

        latest_run = conn.execute(
            """
            SELECT id, source_file, source_year, source_mode, sheet_count, row_count, inserted_count, skipped_count, status, updated_at
            FROM qsse_import_runs
            ORDER BY id DESC
            LIMIT 1
            """
        ).fetchone()

    return {
        "filters": {"year": year},
        "totals": {
            "total_records": int((totals["total_records"] if totals else 0) or 0),
            "event_records": int((totals["event_records"] if totals else 0) or 0),
            "indicator_records": int((totals["indicator_records"] if totals else 0) or 0),
            "records_2025": int((totals["records_2025"] if totals else 0) or 0),
            "records_2026": int((totals["records_2026"] if totals else 0) or 0),
        },
        "registers": [{"register_code": row["register_code"], "row_count": int(row["row_count"] or 0)} for row in registers],
        "years": [{"source_year": int(row["source_year"]), "row_count": int(row["row_count"] or 0)} for row in years],
        "latest_live_run": dict(latest_live_run) if latest_live_run else None,
        "latest_run": dict(latest_run) if latest_run else None,
    }


@router.get("/qsse/records")
def qsse_records(
    year: Optional[int] = Query(None),
    include_backlog: bool = Query(True),
    register_code: Optional[str] = Query(None),
    search: Optional[str] = Query(None),
    limit: int = Query(200, ge=1, le=5000),
):
    where_clauses = ["1 = 1"]
    params: list[Any] = []

    if year is not None:
        where_clauses.append("source_year = ?")
        params.append(int(year))

    normalized_register = (register_code or "").strip().upper()
    if normalized_register and normalized_register != "ALL":
        where_clauses.append("register_code = ?")
        params.append(normalized_register)

    normalized_search = (search or "").strip().lower()
    if normalized_search:
        where_clauses.append(
            "LOWER(COALESCE(title, '') || ' ' || COALESCE(description, '') || ' ' || COALESCE(site, '') || ' ' || COALESCE(agency, '')) LIKE ?"
        )
        params.append(f"%{normalized_search}%")

    where_sql = " AND ".join(where_clauses)

    with connect_qsse_db() as conn:
        total = conn.execute(
            f"SELECT COUNT(*) AS n FROM qsse_records WHERE {where_sql}",
            params,
        ).fetchone()["n"]

        rows = conn.execute(
            f"""
            SELECT
                id,
                source_file,
                source_year,
                source_mode,
                sheet_name,
                sheet_kind,
                row_index,
                register_code,
                record_kind,
                agency,
                entity,
                person,
                site,
                theme,
                title,
                description,
                cause,
                treatment,
                corrective_action,
                action_label,
                pilot,
                status,
                severity,
                date_event,
                date_closed,
                date_saisie,
                amount_value,
                document_reference,
                (
                    SELECT COUNT(*)
                    FROM qsse_documents doc
                    WHERE doc.qsse_record_id = qsse_records.id
                ) AS attachment_count,
                metrics_json,
                updated_at
            FROM qsse_records
            WHERE {where_sql}
            ORDER BY source_year DESC, COALESCE(date_event, '') DESC, id DESC
            LIMIT ?
            """,
            [*params, int(limit)],
        ).fetchall()

    return {
        "total": int(total or 0),
        "limit": int(limit),
        "items": [dict(row) for row in rows],
    }


@router.get("/qsse/records/{record_id}/documents")
def qsse_record_documents(record_id: int):
    with connect_qsse_db() as conn:
        record = _qsse_record_row(conn, record_id)
        if record is None:
            raise HTTPException(status_code=404, detail="Enregistrement QSSE introuvable.")

        rows = conn.execute(
            """
            SELECT id, stored_name, original_name, content_type, file_size, created_at
            FROM qsse_documents
            WHERE qsse_record_id = ?
            ORDER BY created_at DESC, id DESC
            """,
            (int(record_id),),
        ).fetchall()

    return {
        "record_id": int(record_id),
        "documents": [_serialize_qsse_document(record_id, row) for row in rows],
    }


@router.patch("/qsse/records/{record_id}")
def update_qsse_record_cell(record_id: int, data: QsseRecordCellUpdateSchema):
    mode = str(data.mode or "column").strip().lower()
    field = str(data.field or "").strip()
    if not field:
        raise HTTPException(status_code=400, detail="Champ QSSE invalide.")

    with connect_qsse_db() as conn:
        record = conn.execute(
            """
            SELECT id, record_kind, register_code, metrics_json
            FROM qsse_records
            WHERE id = ?
            """,
            (int(record_id),),
        ).fetchone()
        if record is None:
            raise HTTPException(status_code=404, detail="Enregistrement QSSE introuvable.")

        if mode == "column":
            if field not in QSSE_EDITABLE_COLUMN_FIELDS:
                raise HTTPException(status_code=400, detail="Champ QSSE non modifiable dans la table.")

            if field == "amount_value":
                numeric_value = _qsse_number_value(data.value)
                amount_text = "" if numeric_value is None else str(numeric_value)
                conn.execute(
                    """
                    UPDATE qsse_records
                    SET amount_value = ?, amount_text = ?, updated_at = datetime('now')
                    WHERE id = ?
                    """,
                    (numeric_value, amount_text, int(record_id)),
                )
            else:
                conn.execute(
                    f"UPDATE qsse_records SET {field} = ?, updated_at = datetime('now') WHERE id = ?",
                    (_qsse_text_value(data.value), int(record_id)),
                )

        elif mode == "metric":
            metric_key = field.lower()
            if not metric_key or len(metric_key) > 120:
                raise HTTPException(status_code=400, detail="Clé métrique QSSE invalide.")

            aliases = [str(alias or "").strip().lower() for alias in (data.aliases or [])]
            aliases = [alias for alias in aliases if alias and len(alias) <= 120]
            if metric_key not in aliases:
                aliases.append(metric_key)

            metrics = _qsse_json_object(record["metrics_json"])
            for alias in aliases:
                metrics.pop(alias, None)

            metric_value = _qsse_text_value(data.value)
            if metric_value:
                metrics[metric_key] = metric_value

            conn.execute(
                """
                UPDATE qsse_records
                SET metrics_json = ?, updated_at = datetime('now')
                WHERE id = ?
                """,
                (json.dumps(metrics, ensure_ascii=False), int(record_id)),
            )
        else:
            raise HTTPException(status_code=400, detail="Mode de mise à jour QSSE invalide.")

        conn.execute(
            """
            UPDATE qsse_rex_drafts
            SET status = 'stale', updated_at = datetime('now')
            WHERE qsse_record_id = ?
            """,
            (int(record_id),),
        )

        updated = conn.execute(
            """
            SELECT id, updated_at
            FROM qsse_records
            WHERE id = ?
            """,
            (int(record_id),),
        ).fetchone()

    return {
        "ok": True,
        "record_id": int(record_id),
        "mode": mode,
        "field": field,
        "updated_at": updated["updated_at"] if updated else "",
    }


@router.post("/qsse/records/{record_id}/documents")
async def upload_qsse_record_document(record_id: int, file: UploadFile = File(...)):
    original_name = Path(file.filename or "document").name
    suffix = Path(original_name).suffix.lower()
    if suffix not in QSSE_ATTACHMENT_EXTENSIONS:
        raise HTTPException(status_code=400, detail="Format non supporté. Utilisez uniquement un PDF.")

    content = await file.read()
    if not content:
        raise HTTPException(status_code=400, detail="Fichier vide.")
    if len(content) > QSSE_ATTACHMENT_SIZE_LIMIT:
        raise HTTPException(status_code=400, detail="Fichier trop volumineux (max 25 Mo).")

    with connect_qsse_db() as conn:
        record = _qsse_record_row(conn, record_id)
        if record is None:
            raise HTTPException(status_code=404, detail="Enregistrement QSSE introuvable.")
        if str(record["record_kind"] or "").lower() != "event" or str(record["register_code"] or "").upper() not in QSSE_ATTACHMENT_REGISTER_CODES:
            raise HTTPException(status_code=400, detail="Les annexes QSSE sont activées pour les enregistrements FNC, PASD, BP et FAE.")

        target_dir = QSSE_ATTACHMENTS_ROOT / str(int(record_id))
        target_dir.mkdir(parents=True, exist_ok=True)
        stored_name = f"{secrets.token_hex(12)}{suffix}"
        target_path = target_dir / stored_name
        target_path.write_bytes(content)

        try:
            cursor = conn.execute(
                """
                INSERT INTO qsse_documents (
                    qsse_record_id,
                    stored_name,
                    original_name,
                    content_type,
                    file_size
                ) VALUES (?, ?, ?, ?, ?)
                """,
                (
                    int(record_id),
                    stored_name,
                    original_name,
                    file.content_type or "application/octet-stream",
                    len(content),
                ),
            )
            row = conn.execute(
                """
                SELECT id, stored_name, original_name, content_type, file_size, created_at
                FROM qsse_documents
                WHERE id = ?
                """,
                (int(cursor.lastrowid),),
            ).fetchone()
        except Exception:
            if target_path.exists():
                target_path.unlink(missing_ok=True)
            raise

    return {
        "ok": True,
        "record_id": int(record_id),
        "document": _serialize_qsse_document(record_id, row),
    }


@router.delete("/qsse/documents/{document_id}")
def delete_qsse_document(document_id: int):
    with connect_qsse_db() as conn:
        row = conn.execute(
            """
            SELECT id, qsse_record_id, stored_name
            FROM qsse_documents
            WHERE id = ?
            """,
            (int(document_id),),
        ).fetchone()
        if row is None:
            raise HTTPException(status_code=404, detail="Document QSSE introuvable.")

        target_path = QSSE_ATTACHMENTS_ROOT / str(int(row["qsse_record_id"])) / str(row["stored_name"] or "")
        conn.execute("DELETE FROM qsse_documents WHERE id = ?", (int(document_id),))

    if target_path.exists():
        target_path.unlink(missing_ok=True)

    return {
        "ok": True,
        "document_id": int(document_id),
        "record_id": int(row["qsse_record_id"]),
    }


@router.get("/qsse/fnc-analysis/presentations")
def list_qsse_fnc_analysis_presentations(year: Optional[int] = Query(None)):
    params: list[Any] = [QSSE_ANALYSIS_FNC_CODE]
    where_sql = "analysis_code = ?"
    if year is not None:
        where_sql += " AND source_year = ?"
        params.append(int(year))

    with connect_qsse_db() as conn:
        rows = conn.execute(
            f"""
            SELECT id, analysis_code, source_year, stored_name, original_name, content_type, file_size, created_at
            FROM qsse_analysis_documents
            WHERE {where_sql}
            ORDER BY source_year DESC, created_at DESC, id DESC
            """,
            params,
        ).fetchall()

    return {
        "analysis_code": QSSE_ANALYSIS_FNC_CODE,
        "year": int(year) if year is not None else None,
        "documents": [_serialize_qsse_analysis_document(row) for row in rows],
    }


@router.post("/qsse/fnc-analysis/presentations")
async def upload_qsse_fnc_analysis_presentation(
    file: UploadFile = File(...),
    year: Optional[int] = Query(None),
):
    original_name = Path(file.filename or "presentation.pptx").name
    suffix = Path(original_name).suffix.lower()
    if suffix not in QSSE_ANALYSIS_PPTX_EXTENSIONS:
        raise HTTPException(status_code=400, detail="Format non supporté. Utilisez un fichier PPTX.")

    content = await file.read()
    if not content:
        raise HTTPException(status_code=400, detail="Fichier vide.")
    if len(content) > QSSE_ATTACHMENT_SIZE_LIMIT:
        raise HTTPException(status_code=400, detail="Fichier trop volumineux (max 25 Mo).")

    source_year = int(year or 0)
    target_dir = QSSE_ANALYSIS_ATTACHMENTS_ROOT / QSSE_ANALYSIS_FNC_CODE / str(source_year)
    target_dir.mkdir(parents=True, exist_ok=True)
    stored_name = f"{secrets.token_hex(12)}{suffix}"
    target_path = target_dir / stored_name
    target_path.write_bytes(content)

    with connect_qsse_db() as conn:
        try:
            cursor = conn.execute(
                """
                INSERT INTO qsse_analysis_documents (
                    analysis_code,
                    source_year,
                    stored_name,
                    original_name,
                    content_type,
                    file_size
                ) VALUES (?, ?, ?, ?, ?, ?)
                """,
                (
                    QSSE_ANALYSIS_FNC_CODE,
                    source_year,
                    stored_name,
                    original_name,
                    file.content_type or "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    len(content),
                ),
            )
            row = conn.execute(
                """
                SELECT id, analysis_code, source_year, stored_name, original_name, content_type, file_size, created_at
                FROM qsse_analysis_documents
                WHERE id = ?
                """,
                (int(cursor.lastrowid),),
            ).fetchone()
        except Exception:
            if target_path.exists():
                target_path.unlink(missing_ok=True)
            raise

    return {
        "ok": True,
        "analysis_code": QSSE_ANALYSIS_FNC_CODE,
        "document": _serialize_qsse_analysis_document(row),
    }


@router.delete("/qsse/fnc-analysis/presentations/{document_id}")
def delete_qsse_fnc_analysis_presentation(document_id: int):
    with connect_qsse_db() as conn:
        row = conn.execute(
            """
            SELECT id, analysis_code, source_year, stored_name
            FROM qsse_analysis_documents
            WHERE id = ?
            """,
            (int(document_id),),
        ).fetchone()
        if row is None:
            raise HTTPException(status_code=404, detail="Présentation FNC introuvable.")

        target_path = (
            QSSE_ANALYSIS_ATTACHMENTS_ROOT
            / str(row["analysis_code"] or QSSE_ANALYSIS_FNC_CODE)
            / str(int(row["source_year"] or 0))
            / str(row["stored_name"] or "")
        )
        conn.execute("DELETE FROM qsse_analysis_documents WHERE id = ?", (int(document_id),))

    if target_path.exists():
        target_path.unlink(missing_ok=True)

    return {
        "ok": True,
        "document_id": int(document_id),
    }


@router.post("/qsse/fnc-analysis/presentations/generate")
def generate_qsse_fnc_analysis_presentation(
    year: Optional[int] = Query(None),
    template_mode: str = Query("exploitation"),
):
    source_year = int(year or 0)
    mode = str(template_mode or "exploitation").strip().lower()
    if mode not in QSSE_ANALYSIS_TEMPLATE_MODES:
        raise HTTPException(status_code=400, detail="Mode de template invalide. Utilisez codir ou exploitation.")

    with connect_qsse_db() as conn:
        if year is None:
            rows = conn.execute(
                """
                SELECT id, title, agency, status, severity, amount_value, date_closed
                FROM qsse_records
                WHERE register_code = 'FNC' AND record_kind = 'event'
                ORDER BY COALESCE(date_event, '') DESC, id DESC
                """
            ).fetchall()
        else:
            rows = conn.execute(
                """
                SELECT id, title, agency, status, severity, amount_value, date_closed
                FROM qsse_records
                WHERE register_code = 'FNC' AND record_kind = 'event' AND source_year = ?
                ORDER BY COALESCE(date_event, '') DESC, id DESC
                """,
                (source_year,),
            ).fetchall()

    if not rows:
        raise HTTPException(status_code=400, detail="Aucune FNC disponible pour générer la présentation.")

    content = _build_fnc_analysis_pptx(rows, year, mode)
    target_dir = QSSE_ANALYSIS_ATTACHMENTS_ROOT / QSSE_ANALYSIS_FNC_CODE / str(source_year)
    target_dir.mkdir(parents=True, exist_ok=True)
    stored_name = f"{secrets.token_hex(12)}.pptx"
    target_path = target_dir / stored_name
    target_path.write_bytes(content)

    original_name = (
        f"FNC_Analyse_{mode.upper()}_{source_year if source_year else 'ALL'}_{datetime.now().strftime('%Y%m%d_%H%M')}.pptx"
    )

    with connect_qsse_db() as conn:
        try:
            cursor = conn.execute(
                """
                INSERT INTO qsse_analysis_documents (
                    analysis_code,
                    source_year,
                    stored_name,
                    original_name,
                    content_type,
                    file_size
                ) VALUES (?, ?, ?, ?, ?, ?)
                """,
                (
                    QSSE_ANALYSIS_FNC_CODE,
                    source_year,
                    stored_name,
                    original_name,
                    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    len(content),
                ),
            )
            row = conn.execute(
                """
                SELECT id, analysis_code, source_year, stored_name, original_name, content_type, file_size, created_at
                FROM qsse_analysis_documents
                WHERE id = ?
                """,
                (int(cursor.lastrowid),),
            ).fetchone()
        except Exception:
            if target_path.exists():
                target_path.unlink(missing_ok=True)
            raise

    return {
        "ok": True,
        "analysis_code": QSSE_ANALYSIS_FNC_CODE,
        "template_mode": mode,
        "document": _serialize_qsse_analysis_document(row),
    }


@router.get("/qsse/analysis-stats")
def qsse_analysis_stats(
    year: Optional[int] = Query(None),
    include_backlog: bool = Query(True),
    backlog_year: Optional[int] = Query(None),
):
    """Retorna statistiques QSSE avec données trimestrales et comparaison année N-1"""

    def get_date_year(date_str: Any) -> Optional[int]:
        raw = str(date_str or "").strip()
        if not raw or raw == "/":
            return None
        if len(raw) >= 10 and raw[4] == "-" and raw[7] == "-":
            try:
                return int(raw[:4])
            except ValueError:
                return None
        if len(raw) == 10 and raw[2] == "/" and raw[5] == "/":
            try:
                return int(raw[6:10])
            except ValueError:
                return None
        try:
            return datetime.fromisoformat(raw.replace("Z", "+00:00")).year
        except ValueError:
            return None
    
    def get_quarter_from_date(date_str: str) -> Optional[str]:
        """Retourne 'Q1', 'Q2', 'Q3', 'Q4' à partir d'une date YYYY-MM-DD"""
        if not date_str or len(date_str) < 7:
            return None
        month = int(date_str[5:7])
        quarter = (month - 1) // 3 + 1
        year_part = date_str[:4]
        return f"{year_part}-Q{quarter}"
    
    def fetch_year_data(
        conn,
        target_year: Optional[int],
        *,
        include_backlog_scope: bool,
        backlog_year_scope: Optional[int],
    ) -> dict:
        """Récupère toutes les données QSSE pour une année donnée"""
        yc = "AND source_year = ?" if target_year else ""
        ya: list[Any] = [target_year] if target_year else []
        reference_year = int(target_year or backlog_year_scope or datetime.now().year)

        def is_backlog_record(row: Any) -> bool:
            if str(row["register_code"] or "").strip().upper() != "FNC":
                return False
            try:
                source_year = int(row["source_year"] or 0)
            except (TypeError, ValueError):
                return False
            closed_year = get_date_year(row["date_closed"])
            return source_year == reference_year and closed_year == reference_year - 1

        # FNC with date_event for quarters
        fnc_rows = conn.execute(
            f"SELECT id, agency, source_year, register_code, date_closed, amount_value, amount_text, date_event FROM qsse_records WHERE register_code='FNC' AND record_kind='event' {yc} ORDER BY id DESC",
            ya,
        ).fetchall()
        if not include_backlog_scope:
            fnc_rows = [row for row in fnc_rows if not is_backlog_record(row)]
        fnc_total = len(fnc_rows)
        fnc_closed = sum(1 for r in fnc_rows if str(r["date_closed"] or "").strip())
        fnc_cost = sum(float(r["amount_value"] or 0.0) for r in fnc_rows)
        fnc_with_cost_analysis = sum(
            1
            for r in fnc_rows
            if r["amount_value"] is not None or str(r["amount_text"] or "").strip()
        )
        fnc_without_cost_analysis = max(0, fnc_total - fnc_with_cost_analysis)
        fnc_by_ag: dict[str, dict] = {}
        fnc_by_quarter: Counter[str] = Counter()
        
        for r in fnc_rows:
            ag = str(r["agency"] or "").strip() or "—"
            d = fnc_by_ag.setdefault(ag, {"agency": ag, "total": 0, "closed": 0, "open": 0, "cost": 0.0})
            d["total"] += 1
            if str(r["date_closed"] or "").strip():
                d["closed"] += 1
            else:
                d["open"] += 1
            d["cost"] += float(r["amount_value"] or 0.0)
            # Quarterly breakdown
            de = str(r["date_event"] or "").strip()
            q = get_quarter_from_date(de)
            if q:
                fnc_by_quarter[q] += 1
        
        fnc_by_agency = sorted(fnc_by_ag.values(), key=lambda x: -x["total"])
        for item in fnc_by_agency:
            item["cost"] = int(round(item["cost"]))

        # PASD
        pasd_rows = conn.execute(
            f"SELECT agency, date_event FROM qsse_records WHERE register_code='PASD' AND record_kind='event' {yc}",
            ya,
        ).fetchall()
        pasd_ag: Counter[str] = Counter()
        pasd_month: Counter[str] = Counter()
        pasd_by_quarter: Counter[str] = Counter()
        
        for r in pasd_rows:
            pasd_ag[str(r["agency"] or "").strip() or "—"] += 1
            de = str(r["date_event"] or "").strip()
            if len(de) >= 7:
                pasd_month[de[:7]] += 1
            q = get_quarter_from_date(de)
            if q:
                pasd_by_quarter[q] += 1

        # AT
        at_rows = conn.execute(
            f"SELECT agency, date_event FROM qsse_records WHERE register_code='AT' AND record_kind='event' {yc}",
            ya,
        ).fetchall()
        at_ag: Counter[str] = Counter()
        at_by_quarter: Counter[str] = Counter()
        
        for r in at_rows:
            at_ag[str(r["agency"] or "").strip() or "—"] += 1
            de = str(r["date_event"] or "").strip()
            q = get_quarter_from_date(de)
            if q:
                at_by_quarter[q] += 1

        # BP
        bp_rows = conn.execute(
            f"SELECT agency, date_event FROM qsse_records WHERE register_code='BP' AND record_kind='event' {yc}",
            ya,
        ).fetchall()
        bp_ag: Counter[str] = Counter()
        bp_by_quarter: Counter[str] = Counter()
        
        for r in bp_rows:
            bp_ag[str(r["agency"] or "").strip() or "—"] += 1
            de = str(r["date_event"] or "").strip()
            q = get_quarter_from_date(de)
            if q:
                bp_by_quarter[q] += 1

        # FAE
        fae_rows = conn.execute(
            f"SELECT agency, date_event FROM qsse_records WHERE register_code='FAE' AND record_kind='event' {yc}",
            ya,
        ).fetchall()
        fae_ag: Counter[str] = Counter()
        fae_by_quarter: Counter[str] = Counter()
        
        for r in fae_rows:
            fae_ag[str(r["agency"] or "").strip() or "—"] += 1
            de = str(r["date_event"] or "").strip()
            q = get_quarter_from_date(de)
            if q:
                fae_by_quarter[q] += 1

        # Combined per-agency table
        all_agencies = set(fnc_by_ag) | set(pasd_ag) | set(at_ag) | set(bp_ag) | set(fae_ag)
        fnc_map = {item["agency"]: item for item in fnc_by_agency}
        by_agency_all = sorted(
            [
                {
                    "agency": ag,
                    "fnc": fnc_map.get(ag, {}).get("total", 0),
                    "fnc_open": fnc_map.get(ag, {}).get("open", 0),
                    "pasd": pasd_ag.get(ag, 0),
                    "at": at_ag.get(ag, 0),
                    "bp": bp_ag.get(ag, 0),
                    "fae": fae_ag.get(ag, 0),
                }
                for ag in all_agencies
            ],
            key=lambda x: -(x["fnc"] + x["pasd"] + x["at"] + x["bp"] + x["fae"]),
        )

        return {
            "fnc": {
                "total": fnc_total,
                "closed": fnc_closed,
                "open": fnc_total - fnc_closed,
                "cost_total": int(round(fnc_cost)),
                "with_cost_analysis": int(fnc_with_cost_analysis),
                "without_cost_analysis": int(fnc_without_cost_analysis),
                "by_agency": fnc_by_agency,
                "by_quarter": [{"quarter": k, "total": v} for k, v in sorted(fnc_by_quarter.items())],
            },
            "pasd": {
                "total": len(pasd_rows),
                "by_agency": [{"agency": k, "total": v} for k, v in pasd_ag.most_common()],
                "by_month": [{"month": k, "total": v} for k, v in sorted(pasd_month.items())],
                "by_quarter": [{"quarter": k, "total": v} for k, v in sorted(pasd_by_quarter.items())],
            },
            "at": {
                "total": len(at_rows),
                "by_agency": [{"agency": k, "total": v} for k, v in at_ag.most_common()],
                "by_quarter": [{"quarter": k, "total": v} for k, v in sorted(at_by_quarter.items())],
            },
            "bp": {
                "total": len(bp_rows),
                "by_agency": [{"agency": k, "total": v} for k, v in bp_ag.most_common()],
                "by_quarter": [{"quarter": k, "total": v} for k, v in sorted(bp_by_quarter.items())],
            },
            "fae": {
                "total": len(fae_rows),
                "by_agency": [{"agency": k, "total": v} for k, v in fae_ag.most_common()],
                "by_quarter": [{"quarter": k, "total": v} for k, v in sorted(fae_by_quarter.items())],
            },
            "by_agency_all": by_agency_all,
        }

    with connect_qsse_db() as conn:
        current_data = fetch_year_data(
            conn,
            year,
            include_backlog_scope=include_backlog,
            backlog_year_scope=backlog_year,
        )
        
        # Year-over-year: fetch previous year if current year is specified
        prev_year_data = None
        if year:
            prev_year_data = fetch_year_data(
                conn,
                year - 1,
                include_backlog_scope=include_backlog,
                backlog_year_scope=year - 1,
            )

    return {
        "year": year,
        **current_data,
        "year_over_year": {
            "previous_year": year - 1 if year else None,
            "fnc_total_prev": prev_year_data["fnc"]["total"] if prev_year_data else None,
            "pasd_total_prev": prev_year_data["pasd"]["total"] if prev_year_data else None,
            "at_total_prev": prev_year_data["at"]["total"] if prev_year_data else None,
            "bp_total_prev": prev_year_data["bp"]["total"] if prev_year_data else None,
            "fae_total_prev": prev_year_data["fae"]["total"] if prev_year_data else None,
        } if prev_year_data else None,
    }


@router.get("/qsse/records/{record_id}/rex-draft")
def qsse_record_rex_draft(record_id: int):
    try:
        draft = _qsse_rex.get_for_record(record_id)
    except ValueError as exc:
        raise HTTPException(status_code=404, detail=str(exc)) from exc

    return {
        "record_id": int(record_id),
        "draft": draft,
    }


@router.post("/qsse/records/{record_id}/rex-draft/generate")
def generate_qsse_record_rex_draft(record_id: int):
    try:
        draft = _qsse_rex.generate_for_record(record_id)
    except ValueError as exc:
        message = str(exc)
        status_code = 404 if "introuvable" in message.lower() else 400
        raise HTTPException(status_code=status_code, detail=message) from exc

    return {
        "ok": True,
        "record_id": int(record_id),
        "draft": draft,
    }


@router.post("/qsse/refresh-live")
def qsse_refresh_live(replace_existing: bool = Query(True)):
    workbook_path = PROJECT_ROOT / "storage" / "documents" / "Suivi des indicateurs 2026.xlsx"
    service = QsseImportService()

    try:
        result = service.import_sources(
            (
                WorkbookSource(
                    path=workbook_path,
                    source_year=2026,
                    source_mode="live",
                ),
            ),
            replace_existing=replace_existing,
        )
    except FileNotFoundError as exc:
        raise HTTPException(status_code=404, detail=str(exc)) from exc

    return {
        "message": "QSSE 2026 refreshed successfully.",
        "replace_existing": bool(replace_existing),
        "result": result,
    }
